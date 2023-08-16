#############  topo_placer_script_Andrew.py  ###########
#
# Edited by Andrew Gumbert for Thomas Hansen's InfoPos Project
# In Dr. Gina Kuperberg's NeuroCognition of Language Lab
# August 2023
#
# Script is mainly based on original "topo_placer_script.py"
# previously used in the Kuperberg lab.
#
# Credit to Edward Alexander for the original 
# topo_placer_script.py
# edward.alexander@tufts.edu
# This version adapts Alexander's version for the 
# InfoPos project and fixes bugs related to the scale 
# of certain items on the output pptx. Aside from these 
# changes, most of the code was created by Alexander. 
#
# Expects topo maps from "Voltage_Map_Generator_Andrew.m"
# to be in topo_plots folder within this directory
#
# Will create an "extra_images" folder in topo_maps to store
# scaled and cropped topo maps.
#
# Uses pptx with dimensions of "widescreen.pptx" to populate
# a new pptx called "Export_all_topoplots.pptx". This new
# slideshow contains orderly rows of all topo plots from
# the "extra_images" folder.

import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from pptx.dml.color import ColorFormat, RGBColor
import os

ppt_name                       = 'Export_all_topoplots'

# for Tom's computer:
source_image_directory  = "S:\PROJECTS\InfoPos\topo_plots"
ppt_output_directory    = "S:\PROJECTS\InfoPos"

# for Andrew's computer (in S-drive):
#source_image_directory = "/Volumes/as_rsch_NCL02$/PROJECTS/InfoPos/topo_plots"
#ppt_output_directory = "/Volumes/as_rsch_NCL02$/PROJECTS/InfoPos/"


# for Andrew's computer (outside of S-drive):
#source_image_directory         = '/Users/Andrew/Desktop/MatLab/Andrew_PreProcess_FileStructure/topo_plots'  # name of folder script draws images from
#ppt_output_directory           = '/Users/Andrew/Desktop/MatLab/Andrew_PreProcess_FileStructure'

#-----------------------------#
#- Set up ppt                -#
#-----------------------------#
# No way to programmatically change aspect ratio, so save a blank wide pptx and start with that
prs = Presentation('widescreen.pptx')
#prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "ERP TopoPlots"
subtitle.text = "(to check over)"

# Change to suit
conditions = [13, 14, 15, 16, 17, 18]
condition_names = ["1st_info Minus\n1st_uninfo_Diff",
"3rd_info Minus\n3rd_uninfo_Diff",
"Info_exp Minus\nInfo_unexp",
"Uninfo_exp Minus\nUninfo_unexp",
"Schema_1_info Minus\nSchema_1_uninfo",
"Schema_2_info Minus\nSchema_2_uninfo"]
 
#-----------------------------#
#- Crop images               -#
#-----------------------------#

if not os.path.exists(os.path.join(source_image_directory, "extra_images")):
    os.makedirs(os.path.join(source_image_directory, "extra_images"))

#- get just the scale -#
to_crop = Image.open(os.path.join(source_image_directory, str(conditions[0]) + "_" + str(200) + ".gif"))
cropped = to_crop.crop((0, 50, 27, 50+632))
cropped.save(os.path.join(source_image_directory, "extra_images", "scale.gif"))

#- resize topo maps and remove horizontal padding -#

for cond in conditions:
    for time in range(200, 1000, 50):    
        to_crop = Image.open(os.path.join(source_image_directory, str(cond) + "_" + str(time) + ".gif"))
        resized = to_crop.resize((153 + 808, 755)) # this line added by Andrew because otherwise most of map gets cropped out
        cropped = resized.crop((153, 0, 808, 755)) #L,T,R,B
        cropped.save(os.path.join(source_image_directory, "extra_images", str(cond) + "_" + str(time) + ".gif"))

#-----------------------------#
#- Populating ppt            -#
#-----------------------------#

# add a new slide to display topo maps 

slide = prs.slides.add_slide(prs.slide_layouts[6])

# add condition name labels to each row of topo maps

for count, cond in enumerate(conditions):
    txBox = slide.shapes.add_textbox(Inches(0), Inches(((count+1)*1.5)-0.43), Inches(1), Inches(1)) #(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = condition_names[count]
    p.font.bold = False
    p.font.color.rgb = RGBColor(111,111,111)

# add each topo map and its time label

pic = slide.shapes.add_picture(os.path.join(source_image_directory, "extra_images", "scale.gif"), Inches(3.95), Inches(1.35), height=Inches(5.90))

for time in range(200, 600, 50):
    for count, cond in enumerate(conditions):
        top = Inches((count+1)*1.5)
        left = Inches(2.16*((time-200+200)/100))
        height = Inches(1.25)
        pic = slide.shapes.add_picture(os.path.join(source_image_directory, "extra_images", str(cond) + "_" + str(time) + ".gif"), left, top, height=height)

        txBox = slide.shapes.add_textbox(left, Inches(((count+1)*1.5)-0.63), Inches(1), Inches(1)) #(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = str(time) + "-" + str(time+50)
        p.font.size = Pt(13) # this line added by Andrew to make time labels more readable
        p.font.bold = False
        p.font.color.rgb = RGBColor(111,111,111)

# add a second slide to display topo maps

slide = prs.slides.add_slide(prs.slide_layouts[6])

for count, cond in enumerate(conditions):
    txBox = slide.shapes.add_textbox(Inches(0), Inches(((count+1)*1.5)-0.43), Inches(1), Inches(1)) #(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = condition_names[count]
    p.font.bold = False
    p.font.color.rgb = RGBColor(111,111,111)

# add each topo map and its time label

pic = slide.shapes.add_picture(os.path.join(source_image_directory, "extra_images", "scale.gif"), Inches(3.95), Inches(1.35), height=Inches(5.90))

for time in range(600, 1000, 50):
    for count, cond in enumerate(conditions):
        top = Inches((count+1)*1.5)
        left = Inches(2.16*((time-600+200)/100))
        height = Inches(1.25)
        pic = slide.shapes.add_picture(os.path.join(source_image_directory, "extra_images", str(cond) + "_" + str(time) + ".gif"), left, top, height=height)

        txBox = slide.shapes.add_textbox(left, Inches(((count+1)*1.5)-0.63), Inches(1), Inches(1)) #(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = str(time) + "-" + str(time+50)
        p.font.size = Pt(13) # this line added by Andrew to make time labels more readable
        p.font.bold = False
        p.font.color.rgb = RGBColor(111,111,111)

prs.save(os.path.join(ppt_output_directory, ppt_name + '.pptx'))

print("Done.")
